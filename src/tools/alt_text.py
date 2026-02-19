"""Read and write alt text on images in .docx files.

Uses raw XML manipulation via lxml because python-docx doesn't expose
alt text through its API.

Targets:
- 1.1.1: Alt text for non-decorative images; empty alt for decorative
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field

from docx.document import Document
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# XML namespaces
_NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


@dataclass
class AltTextResult:
    """Result of an alt text operation."""
    success: bool
    changes: list[str] = field(default_factory=list)
    error: str = ""


@dataclass
class ImageAltTextInfo:
    """Alt text info for a single image found in the document."""
    paragraph_index: int
    drawing_index: int
    relationship_id: str
    current_alt_text: str
    has_alt_text: bool


def _find_all_images(doc: Document) -> list[ImageAltTextInfo]:
    """Walk document paragraphs and find all images with their alt text state."""
    images: list[ImageAltTextInfo] = []

    for para_idx, paragraph in enumerate(doc.paragraphs):
        para_elem = paragraph._element
        drawings = para_elem.findall(f".//{{{_NS['w']}}}drawing")

        drawing_counter = 0
        for drawing in drawings:
            for pos_type in ("inline", "anchor"):
                for pos_elem in drawing.findall(f"{{{_NS['wp']}}}{pos_type}"):
                    doc_pr = pos_elem.find(f"{{{_NS['wp']}}}docPr")
                    alt_text = ""
                    if doc_pr is not None:
                        alt_text = doc_pr.get("descr", "")

                    blip = pos_elem.find(f".//{{{_NS['a']}}}blip")
                    rel_id = ""
                    if blip is not None:
                        rel_id = blip.get(f"{{{_NS['r']}}}embed", "")

                    images.append(ImageAltTextInfo(
                        paragraph_index=para_idx,
                        drawing_index=drawing_counter,
                        relationship_id=rel_id,
                        current_alt_text=alt_text,
                        has_alt_text=bool(alt_text),
                    ))
                    drawing_counter += 1

    return images


def get_all_alt_text(doc: Document) -> list[ImageAltTextInfo]:
    """Get alt text info for all images in the document."""
    return _find_all_images(doc)


def set_alt_text(
    doc: Document,
    paragraph_index: int,
    alt_text: str,
    drawing_index: int = 0,
) -> AltTextResult:
    """Set alt text on an image in the document.

    Args:
        doc: python-docx Document.
        paragraph_index: Index of the paragraph containing the image.
        alt_text: The alt text to set. Use empty string for decorative images.
        drawing_index: Index of the drawing within the paragraph (default 0).

    Returns:
        AltTextResult with success/failure and change details.
    """
    try:
        if paragraph_index >= len(doc.paragraphs):
            return AltTextResult(
                success=False,
                error=f"Paragraph index {paragraph_index} out of range (doc has {len(doc.paragraphs)} paragraphs)",
            )

        paragraph = doc.paragraphs[paragraph_index]
        para_elem = paragraph._element
        drawings = para_elem.findall(f".//{{{_NS['w']}}}drawing")

        # Flatten all position elements across drawings
        doc_prs: list[tuple] = []
        for drawing in drawings:
            for pos_type in ("inline", "anchor"):
                for pos_elem in drawing.findall(f"{{{_NS['wp']}}}{pos_type}"):
                    doc_pr = pos_elem.find(f"{{{_NS['wp']}}}docPr")
                    if doc_pr is not None:
                        doc_prs.append((doc_pr, pos_elem))

        if drawing_index >= len(doc_prs):
            return AltTextResult(
                success=False,
                error=f"Drawing index {drawing_index} out of range (paragraph has {len(doc_prs)} images)",
            )

        doc_pr, _ = doc_prs[drawing_index]
        old_alt = doc_pr.get("descr", "")
        doc_pr.set("descr", alt_text)

        if alt_text:
            change = f"p_{paragraph_index} img {drawing_index}: alt text set to {alt_text!r} (was {old_alt!r})"
        else:
            change = f"p_{paragraph_index} img {drawing_index}: marked as decorative (was {old_alt!r})"

        logger.info(change)
        return AltTextResult(success=True, changes=[change])

    except Exception as e:
        return AltTextResult(success=False, error=f"Failed to set alt text: {e}")


def set_decorative(
    doc: Document,
    paragraph_index: int,
    drawing_index: int = 0,
) -> AltTextResult:
    """Mark an image as decorative (empty alt text, WCAG 1.1.1).

    Args:
        doc: python-docx Document.
        paragraph_index: Index of the paragraph containing the image.
        drawing_index: Index of the drawing within the paragraph.

    Returns:
        AltTextResult with success/failure.
    """
    return set_alt_text(doc, paragraph_index, "", drawing_index)


def set_alt_text_by_rel_id(
    doc: Document,
    relationship_id: str,
    alt_text: str,
) -> AltTextResult:
    """Set alt text on an image identified by its relationship ID.

    Useful when working from ImageInfo which stores relationship_id.

    Args:
        doc: python-docx Document.
        relationship_id: The rId from the docx relationships.
        alt_text: The alt text to set.

    Returns:
        AltTextResult with success/failure.
    """
    try:
        for para_idx, paragraph in enumerate(doc.paragraphs):
            para_elem = paragraph._element
            drawings = para_elem.findall(f".//{{{_NS['w']}}}drawing")

            drawing_counter = 0
            for drawing in drawings:
                for pos_type in ("inline", "anchor"):
                    for pos_elem in drawing.findall(f"{{{_NS['wp']}}}{pos_type}"):
                        blip = pos_elem.find(f".//{{{_NS['a']}}}blip")
                        if blip is not None:
                            rel_id = blip.get(f"{{{_NS['r']}}}embed", "")
                            if rel_id == relationship_id:
                                return set_alt_text(doc, para_idx, alt_text, drawing_counter)
                        drawing_counter += 1

        return AltTextResult(
            success=False,
            error=f"Image with relationship ID {relationship_id!r} not found",
        )

    except Exception as e:
        return AltTextResult(success=False, error=f"Failed to set alt text: {e}")


# PPTX namespace
_PPTX_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def set_alt_text_pptx(
    prs,
    slide_index: int,
    shape_index: int,
    alt_text: str,
) -> AltTextResult:
    """Set alt text on a PPTX image shape.

    Args:
        prs: python-pptx Presentation object.
        slide_index: 0-based slide index.
        shape_index: 0-based shape index on the slide.
        alt_text: The alt text to set.

    Returns:
        AltTextResult with success/failure.
    """
    try:
        slides = list(prs.slides)
        if slide_index >= len(slides):
            return AltTextResult(
                success=False,
                error=f"Slide index {slide_index} out of range (pres has {len(slides)} slides)",
            )

        slide = slides[slide_index]
        shapes = list(slide.shapes)
        if shape_index >= len(shapes):
            return AltTextResult(
                success=False,
                error=f"Shape index {shape_index} out of range (slide has {len(shapes)} shapes)",
            )

        shape = shapes[shape_index]

        # Find p:cNvPr in the shape's XML
        cNvPr = shape._element.find(f".//{{{_PPTX_NS['p']}}}cNvPr")
        if cNvPr is None:
            return AltTextResult(
                success=False,
                error=f"No cNvPr element found on shape {shape_index} of slide {slide_index}",
            )

        old_alt = cNvPr.get("descr", "")
        cNvPr.set("descr", alt_text)

        if alt_text:
            change = f"slide {slide_index} shape {shape_index}: alt text set to {alt_text[:80]!r} (was {old_alt!r})"
        else:
            change = f"slide {slide_index} shape {shape_index}: marked as decorative (was {old_alt!r})"

        logger.info(change)
        return AltTextResult(success=True, changes=[change])

    except Exception as e:
        return AltTextResult(success=False, error=f"Failed to set alt text on pptx: {e}")
