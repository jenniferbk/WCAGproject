"""Parse .pptx files into a DocumentModel.

Uses python-pptx for high-level access and lxml for raw XML (alt text).
Slide titles become H2 headings, presentation title becomes H1.
Speaker notes are included as paragraphs with style_name="Notes".
"""

from __future__ import annotations

import io
import logging
import re
import statistics
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.models.document import (
    CellInfo,
    ContentOrderItem,
    ContentType,
    DocumentModel,
    DocumentStats,
    FakeHeadingSignals,
    ImageInfo,
    LinkInfo,
    MetadataInfo,
    ParagraphInfo,
    RunInfo,
    TableInfo,
)
from src.tools.docx_parser import ParseResult

logger = logging.getLogger(__name__)

# Default font size in points
DEFAULT_FONT_SIZE_PT = 18.0  # PPTX defaults tend to be larger

# Fake heading scoring weights (same as docx)
_WEIGHT_ALL_BOLD = 0.30
_WEIGHT_FONT_SIZE = 0.25
_WEIGHT_SHORT = 0.20
_WEIGHT_FOLLOWED_BY = 0.15
_WEIGHT_NOT_IN_TABLE = 0.10

# XML namespaces for PPTX
_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def parse_pptx(filepath: str | Path) -> ParseResult:
    """Parse a .pptx file into a DocumentModel.

    Args:
        filepath: Path to the .pptx file.

    Returns:
        ParseResult with success/failure and the DocumentModel.
    """
    filepath = Path(filepath)
    warnings: list[str] = []

    if not filepath.exists():
        return ParseResult(success=False, error=f"File not found: {filepath}")
    if not filepath.suffix.lower() == ".pptx":
        return ParseResult(success=False, error=f"Not a .pptx file: {filepath}")

    try:
        prs = Presentation(str(filepath))
    except Exception as e:
        return ParseResult(success=False, error=f"Failed to open pptx: {e}")

    try:
        # Extract metadata
        metadata = _extract_metadata(prs)

        paragraphs: list[ParagraphInfo] = []
        tables: list[TableInfo] = []
        images: list[ImageInfo] = []
        links: list[LinkInfo] = []
        content_order: list[ContentOrderItem] = []

        para_counter = 0
        table_counter = 0
        image_counter = 0
        link_counter = 0
        all_font_sizes: list[float] = []

        # Add presentation title as H1 if available
        prs_title = (metadata.title or "").strip()
        if prs_title and prs_title != "PowerPoint Presentation":
            para_id = f"p_{para_counter}"
            paragraphs.append(ParagraphInfo(
                id=para_id,
                text=prs_title,
                style_name="Heading 1",
                heading_level=1,
                runs=[RunInfo(text=prs_title)],
            ))
            content_order.append(ContentOrderItem(
                content_type=ContentType.PARAGRAPH, id=para_id,
            ))
            para_counter += 1

        # Process each slide
        for slide_idx, slide in enumerate(prs.slides):
            # Collect surrounding text for image context
            slide_text_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text_parts.append(shape.text_frame.text)
            slide_text = " ".join(slide_text_parts)[:200]

            # Slide title → H2
            title_shape = slide.shapes.title
            title_text = title_shape.text.strip() if title_shape and title_shape.text else ""
            if title_text:
                para_id = f"p_{para_counter}"
                paragraphs.append(ParagraphInfo(
                    id=para_id,
                    text=title_text,
                    style_name="Heading 2",
                    heading_level=2,
                    runs=[RunInfo(text=title_text)],
                ))
                content_order.append(ContentOrderItem(
                    content_type=ContentType.PARAGRAPH, id=para_id,
                ))
                para_counter += 1

            # Process shapes
            for shape_idx, shape in enumerate(slide.shapes):
                # Skip the title shape (already handled)
                if title_shape and shape.shape_id == title_shape.shape_id:
                    continue

                # Images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_info = _extract_image(
                        shape, slide_idx, shape_idx, image_counter, slide_text,
                    )
                    if img_info is not None:
                        images.append(img_info)
                        # Add a placeholder paragraph for the image
                        para_id = f"p_{para_counter}"
                        paragraphs.append(ParagraphInfo(
                            id=para_id,
                            text="",
                            style_name="Normal",
                            image_ids=[img_info.id],
                        ))
                        content_order.append(ContentOrderItem(
                            content_type=ContentType.PARAGRAPH, id=para_id,
                        ))
                        # Update image with paragraph_id
                        images[-1] = img_info.model_copy(update={"paragraph_id": para_id})
                        para_counter += 1
                        image_counter += 1

                # Tables
                elif shape.has_table:
                    tbl_id = f"tbl_{table_counter}"
                    tbl_info = _parse_table(shape.table, tbl_id)
                    tables.append(tbl_info)
                    content_order.append(ContentOrderItem(
                        content_type=ContentType.TABLE, id=tbl_id,
                    ))
                    table_counter += 1

                # Text frames (not title)
                elif shape.has_text_frame:
                    for pptx_para in shape.text_frame.paragraphs:
                        text = pptx_para.text.strip()
                        if not text:
                            continue

                        para_id = f"p_{para_counter}"
                        runs: list[RunInfo] = []
                        para_links: list[LinkInfo] = []
                        para_font_sizes: list[float] = []

                        for run in pptx_para.runs:
                            if not run.text:
                                continue
                            font_size = run.font.size.pt if run.font.size else None
                            if font_size is not None:
                                para_font_sizes.append(font_size)

                            color_hex = None
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    rgb = run.font.color.rgb
                                    if rgb is not None:
                                        color_hex = f"#{rgb}"
                            except AttributeError:
                                pass

                            runs.append(RunInfo(
                                text=run.text,
                                bold=run.font.bold,
                                italic=run.font.italic,
                                underline=run.font.underline,
                                font_size_pt=font_size,
                                font_name=run.font.name,
                                color=color_hex,
                            ))

                        all_font_sizes.extend(para_font_sizes)

                        # Detect hyperlinks from XML
                        for hl in pptx_para._p.findall(f".//{{{_NS['a']}}}hlinkClick"):
                            r_id = hl.get(f"{{{_NS['r']}}}id", "")
                            url = ""
                            if r_id:
                                try:
                                    rel = slide.part.rels.get(r_id)
                                    if rel:
                                        url = rel._target
                                except Exception:
                                    pass
                            if url:
                                link_id = f"link_{link_counter}"
                                para_links.append(LinkInfo(
                                    id=link_id,
                                    text=text,
                                    url=url,
                                    paragraph_id=para_id,
                                ))
                                links.append(para_links[-1])
                                link_counter += 1

                        paragraphs.append(ParagraphInfo(
                            id=para_id,
                            text=text,
                            style_name="Normal",
                            runs=runs,
                            links=para_links,
                        ))
                        content_order.append(ContentOrderItem(
                            content_type=ContentType.PARAGRAPH, id=para_id,
                        ))
                        para_counter += 1

            # Speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    para_id = f"p_{para_counter}"
                    notes_text = notes_frame.text.strip()
                    paragraphs.append(ParagraphInfo(
                        id=para_id,
                        text=notes_text,
                        style_name="Notes",
                        runs=[RunInfo(text=notes_text)],
                    ))
                    content_order.append(ContentOrderItem(
                        content_type=ContentType.PARAGRAPH, id=para_id,
                    ))
                    para_counter += 1

        # Second pass: score fake headings
        median_font = (
            statistics.median(all_font_sizes) if all_font_sizes
            else DEFAULT_FONT_SIZE_PT
        )
        paragraphs = _score_fake_headings(paragraphs, median_font)

        # Compute stats
        heading_count = sum(
            1 for p in paragraphs if p.heading_level is not None
        )
        images_missing_alt = sum(1 for img in images if not img.alt_text)
        fake_heading_candidates = sum(
            1 for p in paragraphs
            if p.fake_heading_signals is not None and p.fake_heading_signals.score >= 0.5
        )

        stats = DocumentStats(
            paragraph_count=len(paragraphs),
            table_count=len(tables),
            image_count=len(images),
            link_count=len(links),
            heading_count=heading_count,
            images_missing_alt=images_missing_alt,
            fake_heading_candidates=fake_heading_candidates,
        )

        document = DocumentModel(
            source_format="pptx",
            source_path=str(filepath),
            metadata=metadata,
            paragraphs=paragraphs,
            tables=tables,
            images=images,
            links=links,
            content_order=content_order,
            stats=stats,
            parse_warnings=warnings,
        )

        return ParseResult(success=True, document=document, warnings=warnings)

    except Exception as e:
        logger.exception("Error parsing pptx: %s", filepath)
        return ParseResult(success=False, error=f"Parse error: {e}", warnings=warnings)


def _extract_metadata(prs: Presentation) -> MetadataInfo:
    """Extract metadata from PPTX core properties."""
    props = prs.core_properties
    return MetadataInfo(
        title=props.title or "",
        author=props.author or "",
        language=props.language or "",
        subject=props.subject or "",
        created=str(props.created) if props.created else "",
        modified=str(props.modified) if props.modified else "",
    )


def _extract_image(
    shape,
    slide_index: int,
    shape_index: int,
    image_counter: int,
    surrounding_text: str,
) -> ImageInfo | None:
    """Extract image data and alt text from a picture shape."""
    try:
        blob = shape.image.blob
        content_type = shape.image.content_type
    except Exception:
        return None

    # Get alt text from cNvPr[@descr]
    alt_text = ""
    cNvPr = shape._element.find(f".//{{{_NS['p']}}}cNvPr")
    if cNvPr is not None:
        alt_text = cNvPr.get("descr", "")

    # Get dimensions
    width_px = None
    height_px = None
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(blob))
        width_px, height_px = img.size
    except Exception:
        pass

    return ImageInfo(
        id=f"img_{image_counter}",
        image_data=blob,
        content_type=content_type,
        alt_text=alt_text,
        width_px=width_px,
        height_px=height_px,
        surrounding_text=surrounding_text[:200],
        slide_index=slide_index,
        shape_index=shape_index,
    )


def _parse_table(table, tbl_id: str) -> TableInfo:
    """Parse a PPTX table into TableInfo."""
    rows: list[list[CellInfo]] = []

    for row in table.rows:
        cells: list[CellInfo] = []
        for cell in row.cells:
            cell_paras = [p.text for p in cell.text_frame.paragraphs]
            cell_text = "\n".join(cell_paras)

            # Check for merged cells
            grid_span = 1
            if cell.is_merge_origin:
                grid_span = cell.span_width

            cells.append(CellInfo(
                text=cell_text,
                paragraphs=cell_paras,
                grid_span=grid_span,
            ))
        rows.append(cells)

    row_count = len(rows)
    col_count = max((len(r) for r in rows), default=0)

    return TableInfo(
        id=tbl_id,
        rows=rows,
        header_row_count=0,  # PPTX tables don't have header markup
        style_name="",
        row_count=row_count,
        col_count=col_count,
    )


def _score_fake_headings(
    paragraphs: list[ParagraphInfo],
    median_font_size: float,
) -> list[ParagraphInfo]:
    """Score paragraphs for fake heading likelihood.

    Same logic as docx_parser — only scores Normal-style paragraphs.
    """
    result: list[ParagraphInfo] = []

    for i, para in enumerate(paragraphs):
        if para.heading_level is not None or not para.text.strip() or para.is_list_item:
            result.append(para)
            continue

        if para.style_name not in ("Normal", "Body Text", "Body"):
            result.append(para)
            continue

        text_runs = [r for r in para.runs if r.text.strip()]
        if not text_runs:
            result.append(para)
            continue

        all_bold = all(r.bold is True for r in text_runs)

        run_sizes = [r.font_size_pt for r in text_runs if r.font_size_pt is not None]
        max_font_size = max(run_sizes) if run_sizes else None
        font_above_avg = (
            max_font_size is not None and max_font_size >= median_font_size + 2.0
        )

        word_count = len(para.text.split())
        is_short = word_count < 10

        followed_by_non_bold = False
        for j in range(i + 1, min(i + 3, len(paragraphs))):
            next_para = paragraphs[j]
            if next_para.text.strip():
                next_text_runs = [r for r in next_para.runs if r.text.strip()]
                if next_text_runs:
                    followed_by_non_bold = not all(
                        r.bold is True for r in next_text_runs
                    )
                break

        if not all_bold:
            result.append(para)
            continue

        score = (
            _WEIGHT_ALL_BOLD * (1.0 if all_bold else 0.0)
            + _WEIGHT_FONT_SIZE * (1.0 if font_above_avg else 0.0)
            + _WEIGHT_SHORT * (1.0 if is_short else 0.0)
            + _WEIGHT_FOLLOWED_BY * (1.0 if followed_by_non_bold else 0.0)
            + _WEIGHT_NOT_IN_TABLE * 1.0
        )

        signals = FakeHeadingSignals(
            all_runs_bold=all_bold,
            font_size_pt=max_font_size,
            font_size_above_avg=font_above_avg,
            is_short=is_short,
            followed_by_non_bold=followed_by_non_bold,
            not_in_table=True,
            score=round(score, 3),
        )

        updated = para.model_copy(update={"fake_heading_signals": signals})
        result.append(updated)

    return result
