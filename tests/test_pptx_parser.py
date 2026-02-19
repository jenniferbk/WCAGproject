"""Tests for PPTX parser.

Tests against the real test file: Lesson 3 Information Processing.pptx
"""

import pytest
from pathlib import Path

from src.tools.pptx_parser import parse_pptx

TEST_PPTX = Path(__file__).parent.parent / "testdocs" / "Lesson 3 Information Processing.pptx"


@pytest.fixture
def parsed():
    """Parse the test PPTX file once for all tests."""
    result = parse_pptx(str(TEST_PPTX))
    assert result.success, f"Parse failed: {result.error}"
    return result.document


class TestParseResult:
    def test_success(self):
        result = parse_pptx(str(TEST_PPTX))
        assert result.success
        assert result.document is not None

    def test_file_not_found(self):
        result = parse_pptx("/nonexistent/file.pptx")
        assert not result.success
        assert "not found" in result.error.lower()

    def test_wrong_extension(self, tmp_path):
        # Create an actual file with wrong extension so it gets past exists() check
        fake = tmp_path / "test.docx"
        fake.write_text("fake")
        result = parse_pptx(str(fake))
        assert not result.success
        assert "not a .pptx" in result.error.lower()


class TestMetadata:
    def test_format(self, parsed):
        assert parsed.source_format == "pptx"

    def test_title(self, parsed):
        # This file has "PowerPoint Presentation" as default title
        assert parsed.metadata.title == "PowerPoint Presentation"

    def test_author(self, parsed):
        assert parsed.metadata.author == "Amy Ellis"


class TestImages:
    def test_image_count(self, parsed):
        assert parsed.stats.image_count == 14

    def test_images_have_slide_index(self, parsed):
        for img in parsed.images:
            assert img.slide_index is not None
            assert img.shape_index is not None

    def test_images_have_data(self, parsed):
        for img in parsed.images:
            assert img.image_data is not None
            assert len(img.image_data) > 0

    def test_images_have_content_type(self, parsed):
        for img in parsed.images:
            assert img.content_type

    def test_images_have_dimensions(self, parsed):
        for img in parsed.images:
            assert img.width_px is not None and img.width_px > 0
            assert img.height_px is not None and img.height_px > 0

    def test_some_images_have_alt_text(self, parsed):
        with_alt = [img for img in parsed.images if img.alt_text]
        without_alt = [img for img in parsed.images if not img.alt_text]
        # Some have alt text (filenames), some don't
        assert len(with_alt) > 0
        assert len(without_alt) > 0

    def test_images_have_surrounding_text(self, parsed):
        with_context = [img for img in parsed.images if img.surrounding_text]
        assert len(with_context) > 0

    def test_image_ids_sequential(self, parsed):
        for i, img in enumerate(parsed.images):
            assert img.id == f"img_{i}"

    def test_images_have_paragraph_ids(self, parsed):
        for img in parsed.images:
            assert img.paragraph_id.startswith("p_")


class TestTables:
    def test_table_count(self, parsed):
        assert parsed.stats.table_count == 1

    def test_table_has_rows(self, parsed):
        tbl = parsed.tables[0]
        assert tbl.row_count > 0
        assert tbl.col_count > 0

    def test_table_cells_have_text(self, parsed):
        tbl = parsed.tables[0]
        has_text = any(
            cell.text.strip()
            for row in tbl.rows
            for cell in row
        )
        assert has_text


class TestHeadings:
    def test_has_headings(self, parsed):
        assert parsed.stats.heading_count > 0

    def test_slide_titles_are_h2(self, parsed):
        h2s = [p for p in parsed.paragraphs if p.heading_level == 2]
        assert len(h2s) > 0
        # Check that "Schemes" slide title is present
        h2_texts = [p.text for p in h2s]
        assert "Schemes" in h2_texts


class TestSpeakerNotes:
    def test_has_notes(self, parsed):
        notes = [p for p in parsed.paragraphs if p.style_name == "Notes"]
        assert len(notes) > 0

    def test_notes_have_content(self, parsed):
        notes = [p for p in parsed.paragraphs if p.style_name == "Notes"]
        for note in notes:
            assert note.text.strip()


class TestContentOrder:
    def test_content_order_has_entries(self, parsed):
        assert len(parsed.content_order) > 0

    def test_content_order_covers_paragraphs(self, parsed):
        para_ids_in_order = {
            item.id for item in parsed.content_order
            if item.content_type.value == "paragraph"
        }
        para_ids = {p.id for p in parsed.paragraphs}
        assert para_ids == para_ids_in_order

    def test_content_order_covers_tables(self, parsed):
        tbl_ids_in_order = {
            item.id for item in parsed.content_order
            if item.content_type.value == "table"
        }
        tbl_ids = {t.id for t in parsed.tables}
        assert tbl_ids == tbl_ids_in_order


class TestStats:
    def test_paragraph_count(self, parsed):
        assert parsed.stats.paragraph_count == len(parsed.paragraphs)

    def test_images_missing_alt(self, parsed):
        expected = sum(1 for img in parsed.images if not img.alt_text)
        assert parsed.stats.images_missing_alt == expected


class TestAltTextPptx:
    """Test PPTX-specific alt text writing."""

    def test_set_alt_text_pptx(self):
        from pptx import Presentation
        from src.tools.alt_text import set_alt_text_pptx

        prs = Presentation(str(TEST_PPTX))
        # Image on slide 1, shape 1 (Picture 2)
        result = set_alt_text_pptx(prs, 1, 1, "A diagram showing information encoding")
        assert result.success

    def test_set_alt_text_pptx_out_of_range(self):
        from pptx import Presentation
        from src.tools.alt_text import set_alt_text_pptx

        prs = Presentation(str(TEST_PPTX))
        result = set_alt_text_pptx(prs, 999, 0, "test")
        assert not result.success
        assert "out of range" in result.error


class TestMetadataPptx:
    """Test PPTX-specific metadata writing."""

    def test_set_title_pptx(self):
        from pptx import Presentation
        from src.tools.metadata import set_title_pptx

        prs = Presentation(str(TEST_PPTX))
        result = set_title_pptx(prs, "Lesson 3: Information Processing")
        assert result.success
        assert prs.core_properties.title == "Lesson 3: Information Processing"

    def test_set_language_pptx(self):
        from pptx import Presentation
        from src.tools.metadata import set_language_pptx

        prs = Presentation(str(TEST_PPTX))
        result = set_language_pptx(prs, "en")
        assert result.success
        assert prs.core_properties.language == "en"
